import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
from fpdf import FPDF
import tempfile
import os
import time
import requests
import io

# --------------------------------------------------------------------------
# Streamlit 앱: 나만의 AI 학습 사이트 (수정된 전체 파일)
# 주요 변경사항
# - prompt_parts 초기화 오류 수정 (prompt_parts = [])
# - st.status -> st.spinner로 안정화
# - st.secrets에서 API 키 올바르게 읽도록 수정
# - PDF/PPT 업로드 처리에 대한 안정성 및 예외처리 추가
# - 모델 호출 시 prompt를 문자열로 결합
# - 전체를 복붙해서 바로 사용 가능하도록 정리
# --------------------------------------------------------------------------

st.set_page_config(layout="wide", page_title="나만의 AI 학습 사이트")

# API 키 확인 (Streamlit Secrets에 GENAI_API_KEY 키가 있어야 함)
if "GENAI_API_KEY" in st.secrets:
    genai.configure(api_key=st.secrets["GENAI_API_KEY"])
else:
    st.error("설정 오류: Streamlit Secrets에 'GENAI_API_KEY'를 등록해주세요.")
    st.stop()

# 한글 폰트 다운로드 (서버에 폰트가 없으므로 매번 다운로드)
@st.cache_resource
def get_korean_font():
    font_file = "NanumGothic.ttf"
    if not os.path.exists(font_file):
        url = "https://github.com/google/fonts/raw/main/ofl/nanumgothic/NanumGothic-Regular.ttf"
        try:
            response = requests.get(url, timeout=15)
            response.raise_for_status()
            with open(font_file, "wb") as f:
                f.write(response.content)
        except Exception as e:
            # 다운로드 실패 시 경고는 띄우되 앱은 계속 실행
            st.warning(f"한글 폰트 다운로드 실패: {e}")
            return None
    return font_file

FONT_PATH = get_korean_font()

# --------------------------------------------------------------------------
# 기능 함수들
# --------------------------------------------------------------------------

def get_pdf_text(uploaded_file) -> str:
    """업로드된 PDF(UploadedFile 또는 파일 경로)에서 텍스트 추출
    실패 시 에러 문자열을 반환한다.
    """
    try:
        # Streamlit UploadedFile은 file-like이므로 바로 전달 가능
        if hasattr(uploaded_file, "seek"):
            uploaded_file.seek(0)
            reader = PdfReader(uploaded_file)
        else:
            reader = PdfReader(str(uploaded_file))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        return f"PDF 텍스트 추출 실패: {e}"


def get_pptx_text(uploaded_file) -> str:
    """업로드된 PPTX에서 텍스트 추출
    실패 시 에러 문자열을 반환한다.
    """
    try:
        if hasattr(uploaded_file, "seek"):
            uploaded_file.seek(0)
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"PPT 텍스트 추출 실패: {e}"


def upload_media(file):
    """임시 파일로 저장한 뒤 genai에 업로드(비동기 처리 대기 포함).
    실제 genai API 객체의 반환 형태에 따라 약간의 조정이 필요할 수 있음.
    """
    suffix = os.path.splitext(file.name)[1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(file.getvalue())
        tmp_path = tmp.name

    uploaded_file = genai.upload_file(tmp_path)

    # 일부 SDK는 업로드 후 처리 상태를 폴링해야 함
    try:
        # uploaded_file.state.name 형태를 가정하고 폴링
        while getattr(uploaded_file, "state", None) and getattr(uploaded_file.state, "name", "") == "PROCESSING":
            time.sleep(1)
            uploaded_file = genai.get_file(uploaded_file.name)
    except Exception:
        # 폴링이 실패해도 업로드 객체를 반환
        pass

    return uploaded_file


def create_pdf_report(original: str, explanation: str) -> str:
    """간단한 한글 지원 PDF 리포트 생성 후 파일 경로 반환"""
    pdf = FPDF()
    pdf.add_page()
    if FONT_PATH:
        try:
            pdf.add_font('Nanum', '', FONT_PATH, uni=True)
            pdf.set_font('Nanum', size=11)
        except Exception:
            pdf.set_font(size=11)
    else:
        pdf.set_font(size=11)

    # 제목
    pdf.cell(0, 10, txt="AI 학습 리포트", ln=True, align='C')
    pdf.ln(10)

    # 원본 요약
    pdf.set_font(size=10)
    pdf.cell(0, 10, txt="[원본 자료 요약]", ln=True)
    pdf.multi_cell(0, 8, txt=original[:5000] + ("..." if len(original) > 5000 else ""))
    pdf.ln(5)

    # AI 설명
    pdf.cell(0, 10, txt="[AI 상세 설명]", ln=True)
    pdf.multi_cell(0, 8, txt=explanation)

    output_path = tempfile.mktemp(suffix=".pdf")
    pdf.output(output_path)
    return output_path

# --------------------------------------------------------------------------
# 화면 구성
# --------------------------------------------------------------------------

st.title("📚 나만의 영구적인 학습 사이트")
st.markdown("---")
col1, col2 = st.columns(2)

with col1:
    st.header("1. 자료 업로드")
    main_file = st.file_uploader("메인 수업 자료 (PDF)", type=["pdf"]) 
    supp_file = st.file_uploader("보충 자료 (PPT, 영상, 음성 등)", type=["pdf", "pptx", "mp4", "mp3", "wav"]) 

with col2:
    st.header("2. AI 튜터")
    question = st.text_area("질문 또는 요청사항", "이 내용을 바탕으로 핵심 내용을 설명해줘.")

    if st.button("🚀 설명 요청하기"):
        if not main_file:
            st.warning("메인 자료를 먼저 올려주세요.")
        else:
            # 진행 표시
            with st.spinner("AI가 분석 중입니다..."):
                try:
                    # prompt_parts 초기화
                    prompt_parts = []

                    # 1. 메인 자료 텍스트 추가
                    main_text = get_pdf_text(main_file)
                    prompt_parts.append(f"메인 자료 내용:\n{main_text[:30000]}")

                    # 2. 보충 자료 처리
                    if supp_file:
                        ftype = supp_file.name.split('.')[-1].lower()
                        if ftype == 'pdf':
                            supp_text = get_pdf_text(supp_file)
                            prompt_parts.append(f"보충 자료 내용:\n{supp_text[:20000]}")
                        elif ftype in ['pptx', 'ppt']:
                            supp_text = get_pptx_text(supp_file)
                            prompt_parts.append(f"보충 자료 내용:\n{supp_text[:20000]}")
                        else:
                            # 영상/음성
                            st.info("영상/음성 파일 업로드 중... (서버 업로드가 필요합니다)")
                            media = upload_media(supp_file)
                            media_name = getattr(media, 'name', str(media))
                            prompt_parts.append(f"보충 미디어 파일: {media_name}")
                            prompt_parts.append("위 미디어 파일을 참고하세요.")

                    # 3. 질문 추가
                    prompt_parts.append(f"요청 사항: {question}")

                    # 4. 실행 - prompt를 문자열로 합쳐서 보냄
                    st.write("답변 생성 중...")
                    prompt = "\n\n".join(prompt_parts)

                    # 모델 호출 (SDK 버전에 따라 약간의 차이가 있을 수 있음)
                    # 사용 가능한 안정 모델로 변경 (v1beta)
model = genai.GenerativeModel('models/gemini-1.5-pro')
                    response = model.generate_content(prompt)

                    # SDK 반환 형태에 따라 조정 (response.text가 일반적)
                    result_text = getattr(response, 'text', None)
                    if result_text is None:
                        # 문자열로 변환 가능한 경우
                        result_text = str(response)

                    st.session_state['result'] = result_text
                    st.session_state['main_text'] = main_text

                except Exception as e:
                    st.error(f"에러 발생: {e}")

# 결과가 있으면 보여주기
if 'result' in st.session_state:
    st.success("분석 완료!")
    st.write(st.session_state['result'])

    # PDF 다운로드 버튼
    if 'main_text' in st.session_state:
        try:
            pdf_file = create_pdf_report(st.session_state['main_text'], st.session_state['result'])
            with open(pdf_file, "rb") as f:
                st.download_button("📄 PDF로 저장하기", data=f, file_name="study_note.pdf")
        except Exception as e:
            st.error(f"PDF 생성 실패: {e}")
